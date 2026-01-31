"""Document content extraction from various file formats."""
import io
import pdfplumber
from docx import Document
import pandas as pd
from pptx import Presentation
from typing import Optional, Tuple

def extract_text_from_pdf(file_content: bytes, filename: str) -> Tuple[str, bool]:
    """
    Extract text from PDF using pdfplumber.
    Returns: (extracted_text, needs_ocr)
    """
    try:
        pdf_file = io.BytesIO(file_content)
        text_parts = []
        
        with pdfplumber.open(pdf_file) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    text_parts.append(page_text)
        
        extracted_text = "\n\n".join(text_parts)
        
        # Check if we got meaningful text (more than just whitespace)
        if extracted_text.strip() and len(extracted_text.strip()) > 50:
            return extracted_text, False
        else:
            # PDF appears to be scanned/image-based
            return "", True
            
    except Exception as e:
        # If extraction fails, assume OCR is needed
        return "", True

def extract_text_from_docx(file_content: bytes) -> str:
    """Extract text from DOCX file."""
    try:
        doc_file = io.BytesIO(file_content)
        doc = Document(doc_file)
        paragraphs = [para.text for para in doc.paragraphs]
        return "\n".join(paragraphs)
    except Exception as e:
        raise ValueError(f"Failed to extract text from DOCX: {str(e)}")

def extract_text_from_xlsx(file_content: bytes) -> str:
    """Extract text from XLSX file."""
    try:
        xlsx_file = io.BytesIO(file_content)
        df_dict = pd.read_excel(xlsx_file, sheet_name=None, engine='openpyxl')
        
        text_parts = []
        for sheet_name, df in df_dict.items():
            text_parts.append(f"Sheet: {sheet_name}")
            # Convert DataFrame to string representation
            text_parts.append(df.to_string())
            text_parts.append("")
        
        return "\n".join(text_parts)
    except Exception as e:
        raise ValueError(f"Failed to extract text from XLSX: {str(e)}")

def extract_text_from_csv(file_content: bytes) -> str:
    """Extract text from CSV file."""
    try:
        csv_file = io.BytesIO(file_content)
        df = pd.read_csv(csv_file)
        return df.to_string()
    except Exception as e:
        raise ValueError(f"Failed to extract text from CSV: {str(e)}")

def extract_text_from_pptx(file_content: bytes) -> str:
    """Extract text from PPTX file."""
    try:
        pptx_file = io.BytesIO(file_content)
        prs = Presentation(pptx_file)
        
        text_parts = []
        for slide_num, slide in enumerate(prs.slides, 1):
            text_parts.append(f"Slide {slide_num}:")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text_parts.append(shape.text)
            text_parts.append("")
        
        return "\n".join(text_parts)
    except Exception as e:
        raise ValueError(f"Failed to extract text from PPTX: {str(e)}")

def extract_text_from_file(file_content: bytes, filename: str) -> Tuple[str, bool]:
    """
    Extract text from file based on extension.
    Returns: (extracted_text, needs_ocr)
    """
    filename_lower = filename.lower()
    
    if filename_lower.endswith('.pdf'):
        return extract_text_from_pdf(file_content, filename)
    elif filename_lower.endswith('.docx'):
        return extract_text_from_docx(file_content), False
    elif filename_lower.endswith('.xlsx'):
        return extract_text_from_xlsx(file_content), False
    elif filename_lower.endswith('.csv'):
        return extract_text_from_csv(file_content), False
    elif filename_lower.endswith('.pptx'):
        return extract_text_from_pptx(file_content), False
    else:
        # Unknown format or image - needs OCR
        return "", True
