from __future__ import annotations

import csv
import io
from typing import Dict, List, Tuple
from docx import Document
from PIL import Image
from pypdf import PdfReader
from pptx import Presentation
from openpyxl import load_workbook

from app.core.config import Settings
from app.pipeline.ocr_azure import azure_read_document


SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".xlsx", ".csv", ".pptx", ".png", ".jpg", ".jpeg"}


def _extension(filename: str) -> str:
    dot = filename.lower().rfind(".")
    return filename[dot:] if dot >= 0 else ""


def _extract_pdf(data: bytes) -> str:
    reader = PdfReader(io.BytesIO(data))
    texts = []
    for page in reader.pages:
        page_text = page.extract_text() or ""
        texts.append(page_text)
    return "\n".join(t.strip() for t in texts if t.strip())


def _extract_docx(data: bytes) -> str:
    doc = Document(io.BytesIO(data))
    return "\n".join(p.text for p in doc.paragraphs if p.text)


def _extract_pptx(data: bytes) -> str:
    prs = Presentation(io.BytesIO(data))
    texts: List[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if text:
                    texts.append(text)
    return "\n".join(texts)


def _extract_csv(data: bytes) -> str:
    """Convert CSV rows into evidence sentences for downstream processing."""
    decoded = data.decode("utf-8", errors="ignore")
    reader = csv.reader(io.StringIO(decoded))
    rows = list(reader)
    if not rows:
        return decoded

    # Detect header row
    header = [h.strip().lower() for h in rows[0]]
    sentences = []

    # Try structured extraction if we detect relevant columns
    name_cols = [i for i, h in enumerate(header) if h in ("metric name", "metric", "indicator", "name", "kpi")]
    value_cols = [i for i, h in enumerate(header) if h in ("value", "amount", "quantity", "figure")]
    unit_cols = [i for i, h in enumerate(header) if h in ("unit", "units", "uom")]
    year_cols = [i for i, h in enumerate(header) if h in ("year", "reporting year", "period", "fy")]
    cat_cols = [i for i, h in enumerate(header) if h in ("category", "type", "pillar", "section")]
    note_cols = [i for i, h in enumerate(header) if h in ("notes", "description", "remarks", "comment", "source")]

    if name_cols and value_cols:
        # Structured CSV with identifiable columns
        ni, vi = name_cols[0], value_cols[0]
        ui = unit_cols[0] if unit_cols else None
        yi = year_cols[0] if year_cols else None
        ci = cat_cols[0] if cat_cols else None
        nti = note_cols[0] if note_cols else None

        for row in rows[1:]:
            if len(row) <= max(ni, vi):
                continue
            name = row[ni].strip()
            value = row[vi].strip()
            if not name or not value:
                continue
            unit = row[ui].strip() if ui is not None and len(row) > ui else ""
            year = row[yi].strip() if yi is not None and len(row) > yi else ""
            cat = row[ci].strip() if ci is not None and len(row) > ci else ""
            note = row[nti].strip() if nti is not None and len(row) > nti else ""

            parts = []
            if cat:
                parts.append(f"{cat}:")
            parts.append(f"{name}:")
            parts.append(f"{value}")
            if unit:
                parts.append(unit)
            if year:
                parts.append(f"({year})")
            sentence = " ".join(parts) + "."
            if note:
                sentence += f" {note}"
            sentences.append(sentence)
    else:
        # Unstructured CSV: concatenate rows as text
        for row in rows[1:]:
            line = ", ".join(cell.strip() for cell in row if cell.strip())
            if line:
                sentences.append(line + ".")

    # Also include raw header for context
    header_text = "Data columns: " + ", ".join(rows[0]) + "."
    return header_text + "\n" + "\n".join(sentences)


def _extract_xlsx(data: bytes) -> str:
    workbook = load_workbook(io.BytesIO(data), data_only=True)
    sheet = workbook.active
    output = io.StringIO()
    writer = csv.writer(output)
    for row in sheet.iter_rows(values_only=True):
        writer.writerow([cell if cell is not None else "" for cell in row])
    return output.getvalue()


def _extract_image(data: bytes) -> None:
    Image.open(io.BytesIO(data))


def extract_documents(
    files: List[Tuple[str, bytes, str | None]], settings: Settings
) -> Tuple[Dict[str, str], bool]:
    texts: Dict[str, str] = {}
    ocr_used = False
    for filename, data, content_type in files:
        ext = _extension(filename)
        if ext not in SUPPORTED_EXTENSIONS:
            raise ValueError(f"Unsupported file type: {filename}")

        if ext == ".pdf":
            extracted = _extract_pdf(data)
            if len(extracted.strip()) < 200 and settings.azure_docintel_endpoint and settings.azure_docintel_key:
                extracted = azure_read_document(data, content_type or "application/pdf", settings)
                ocr_used = True
        elif ext == ".docx":
            extracted = _extract_docx(data)
        elif ext == ".pptx":
            extracted = _extract_pptx(data)
        elif ext == ".csv":
            extracted = _extract_csv(data)
        elif ext == ".xlsx":
            extracted = _extract_xlsx(data)
        else:
            _extract_image(data)
            if settings.azure_docintel_endpoint and settings.azure_docintel_key:
                extracted = azure_read_document(data, content_type or "image/png", settings)
                ocr_used = True
            else:
                raise ValueError("OCR not configured for image extraction.")
        texts[filename] = extracted
    return texts, ocr_used
